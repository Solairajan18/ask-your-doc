from pypdf import PdfReader
import docx
import pandas as pd
from pptx import Presentation
import json

def extract_text(file):
    if file.name.endswith(".pdf"):
        reader = PdfReader(file)
        return "\n".join([p.extract_text() for p in reader.pages])

    elif file.name.endswith(".docx"):
        doc = docx.Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    elif file.name.endswith(".xlsx"):
        df = pd.read_excel(file)
        return df.to_string()

    elif file.name.endswith(".pptx"):
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    elif file.name.endswith(".txt"):
        return file.read().decode()

    elif file.name.endswith(".json"):
        return json.dumps(json.load(file), indent=2)

    return ""
